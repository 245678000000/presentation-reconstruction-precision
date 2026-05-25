#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


def load_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def color(value: str | None, default=(0, 0, 0)) -> RGBColor:
    if not value or value == "none":
        return RGBColor(*default)
    value = str(value).strip().lstrip("#")
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    try:
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except Exception:
        return RGBColor(*default)


class Map:
    def __init__(self, w: float, h: float, slide_width: float):
        self.w_px = w
        self.h_px = h
        self.slide_w = Inches(slide_width)
        self.slide_h = int(self.slide_w * h / w)
        self.pt_per_px = slide_width * 72 / w

    def x(self, v): return int(self.slide_w * float(v) / self.w_px)
    def y(self, v): return int(self.slide_h * float(v) / self.h_px)
    def w(self, v): return int(self.slide_w * float(v) / self.w_px)
    def h(self, v): return int(self.slide_h * float(v) / self.h_px)
    def pt(self, v): return Pt(max(1, float(v) * self.pt_per_px))


def set_name(shape, name: str | None) -> None:
    if not name:
        return
    try:
        shape._element.xpath(".//p:cNvPr")[0].set("name", str(name))
    except Exception:
        try:
            shape.name = str(name)
        except Exception:
            pass


def clear_effects(shape, el: Dict[str, Any] | None = None) -> None:
    """Prevent Office theme shadows/glow unless the plan explicitly asks for them."""
    el = el or {}
    effects = el.get("effects") or {}
    if effects.get("shadow") not in {None, "none", False}:
        return
    try:
        style = shape._element.find(qn("p:style"))
        if style is not None:
            effect_ref = style.find(qn("a:effectRef"))
            if effect_ref is not None:
                effect_ref.set("idx", "0")
        sppr = shape._element.find(qn("p:spPr"))
        if sppr is not None:
            for child in list(sppr):
                if child.tag in {qn("a:effectLst"), qn("a:effectDag")}:
                    sppr.remove(child)
            sppr.append(OxmlElement("a:effectLst"))
    except Exception:
        pass


def set_font(run, m: Map, spec: Dict[str, Any]) -> None:
    font = spec.get("font_family", "Microsoft YaHei")
    run.font.name = font
    try:
        rpr = run._r.get_or_add_rPr()
        ea = rpr.find(qn("a:ea"))
        if ea is None:
            ea = OxmlElement("a:ea")
            rpr.append(ea)
        ea.set("typeface", font)
    except Exception:
        pass
    run.font.size = m.pt(spec.get("font_size", 24))
    run.font.bold = int(spec.get("font_weight", 400) or 400) >= 600
    run.font.italic = bool(spec.get("italic", False))
    run.font.color.rgb = color(spec.get("color", "#111111"))


def fill_line(shape, el: Dict[str, Any], m: Map) -> None:
    fill = el.get("fill", "#FFFFFF")
    if fill == "none":
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(fill, (255, 255, 255))
    stroke = el.get("stroke", "none")
    if stroke == "none":
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color(stroke)
        shape.line.width = Pt(max(0.25, float(el.get("stroke_width", 1) or 1) * m.pt_per_px))


def text_frame(shape, m: Map, el: Dict[str, Any]):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = m.w((el.get("margin") or {}).get("left", 0))
    tf.margin_right = m.w((el.get("margin") or {}).get("right", 0))
    tf.margin_top = m.h((el.get("margin") or {}).get("top", 0))
    tf.margin_bottom = m.h((el.get("margin") or {}).get("bottom", 0))
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(el.get("valign", "top"), MSO_ANCHOR.TOP)
    return tf


def add_text(slide, m: Map, el: Dict[str, Any]):
    sh = slide.shapes.add_textbox(m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 1)), m.h(el.get("h", 1)))
    set_name(sh, el.get("shape_name") or el.get("id"))
    tf = text_frame(sh, m, el)
    align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}.get(el.get("align", "left"), PP_ALIGN.LEFT)
    for idx, line in enumerate(str(el.get("text", "")).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        for run in p.runs:
            set_font(run, m, el)


def add_rich(slide, m: Map, el: Dict[str, Any]):
    sh = slide.shapes.add_textbox(m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 1)), m.h(el.get("h", 1)))
    set_name(sh, el.get("shape_name") or el.get("id"))
    tf = text_frame(sh, m, el)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(el.get("align", "left"), PP_ALIGN.LEFT)
    for item in el.get("runs", []):
        run = p.add_run()
        run.text = str(item.get("text", ""))
        set_font(run, m, {**el, **item})


def image_path(asset: Dict[str, Any], assets_dir: Path, plan_dir: Path) -> Path | None:
    file = asset.get("file")
    if not file:
        return None
    candidates = [assets_dir / file, plan_dir / file, Path(file)]
    for path in candidates:
        if path.exists():
            return path
    return None


def add_element(slide, m: Map, el: Dict[str, Any], assets: Dict[str, Dict[str, Any]], assets_dir: Path, plan_dir: Path):
    t = el.get("type")
    sh = None
    if t in {"rect", "round_rect"}:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if t == "round_rect" or float(el.get("rx", 0) or 0) > 0 else MSO_SHAPE.RECTANGLE
        sh = slide.shapes.add_shape(shape_type, m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 0)), m.h(el.get("h", 0)))
        fill_line(sh, el, m)
        clear_effects(sh, el)
    elif t == "circle":
        r = float(el.get("r", 0) or 0)
        cx = float(el.get("cx", el.get("x", 0)) or 0)
        cy = float(el.get("cy", el.get("y", 0)) or 0)
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, m.x(cx - r), m.y(cy - r), m.w(2 * r), m.h(2 * r))
        fill_line(sh, el, m)
        clear_effects(sh, el)
    elif t in {"line", "arrow"}:
        sh = slide.shapes.add_connector(1, m.x(el.get("x1", 0)), m.y(el.get("y1", 0)), m.x(el.get("x2", 0)), m.y(el.get("y2", 0)))
        sh.line.color.rgb = color(el.get("stroke", "#000000"))
        sh.line.width = Pt(max(0.25, float(el.get("stroke_width", 1) or 1) * m.pt_per_px))
        clear_effects(sh, el)
    elif t == "text":
        add_text(slide, m, el)
        return
    elif t == "rich_text":
        add_rich(slide, m, el)
        return
    elif t in {"image", "fallback_image"}:
        asset = assets.get(el.get("asset_id"), {})
        path = image_path(asset, assets_dir, plan_dir)
        if path:
            sh = slide.shapes.add_picture(str(path), m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 1)), m.h(el.get("h", 1)))
            clear_effects(sh, el)
    elif t == "table":
        rows, cols = int(el.get("rows", 1) or 1), int(el.get("cols", 1) or 1)
        sh = slide.shapes.add_table(rows, cols, m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 1)), m.h(el.get("h", 1)))
        for r in range(rows):
            for c in range(cols):
                value = ""
                if r < len(el.get("cell_text", [])) and c < len(el["cell_text"][r]):
                    value = str(el["cell_text"][r][c])
                cell = sh.table.cell(r, c)
                cell.text = value
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        set_font(run, m, el)
    if sh is not None:
        set_name(sh, el.get("shape_name") or el.get("id"))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("deck_plan", type=Path)
    parser.add_argument("--out", type=Path, default=Path("work/presentation_reconstruction_pro/reconstructed.pptx"))
    parser.add_argument("--slide-width", type=float, default=13.333333)
    args = parser.parse_args()

    root = args.deck_plan.resolve().parent
    deck = load_json(args.deck_plan)
    canvas = deck.get("canvas", {"width": 1920, "height": 1080})
    m = Map(canvas["width"], canvas["height"], args.slide_width)
    layer_policy = deck.get("layer_policy") or ["background", "fallback", "decor", "structure", "underlay_recovered", "media", "chart", "icon", "text", "annotation", "debug"]
    layers = {name: idx for idx, name in enumerate(layer_policy)}
    prs = Presentation()
    prs.slide_width = m.slide_w
    prs.slide_height = m.slide_h
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]
    for slide_meta in deck.get("slides", []):
        plan_path = root / slide_meta["plan"]
        plan = load_json(plan_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color(plan.get("canvas", {}).get("background", "#FFFFFF"), (255, 255, 255))
        bg.line.fill.background()
        set_name(bg, f"{slide_meta['id']}.background")
        clear_effects(bg, {"effects": {"shadow": "none"}})
        assets = {a.get("id"): a for a in plan.get("assets", [])}
        assets_dir = root / slide_meta["assets_dir"]
        elements = sorted(plan.get("elements", []), key=lambda e: (layers.get(e.get("layer", "media"), 99), float(e.get("z", 0) or 0)))
        for el in elements:
            add_element(slide, m, el, assets, assets_dir, plan_path.parent)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(args.out)


if __name__ == "__main__":
    main()
